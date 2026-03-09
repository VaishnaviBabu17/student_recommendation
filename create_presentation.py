from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf = txBox.text_frame
tf.text = "Student Recommendation System"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf2 = txBox2.text_frame
tf2.text = "Phase 2: Web Development Evaluation"
p2 = tf2.paragraphs[0]
p2.font.size = Pt(28)
p2.alignment = PP_ALIGN.CENTER

# Slide 2: Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Overview"
content = slide.placeholders[1].text_frame
content.text = "About the System"
for point in ["AI-powered student recommendation platform", "Personalized learning resource suggestions", "Covers multiple engineering domains", "Full-stack web application", "Django backend with modern frontend"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 3: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Technology Stack"
content = slide.placeholders[1].text_frame
content.text = "Backend"
for point in ["Django REST Framework", "SQLite/PostgreSQL Database", "JWT Authentication"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
p = content.add_paragraph()
p.text = "Frontend"
p.level = 0
for point in ["React/HTML/CSS/JavaScript", "Responsive UI Design", "AJAX for API calls"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 4: Backend API Development
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "1. Backend API Development"
content = slide.placeholders[1].text_frame
content.text = "RESTful API Implementation"
for point in ["Designed RESTful endpoints following REST principles", "Implemented CRUD operations for student data", "Created recommendation algorithm endpoints", "Tested all routes using Postman/Thunder Client", "Proper HTTP methods (GET, POST, PUT, DELETE)"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 5: API Endpoints
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "API Endpoints"
content = slide.placeholders[1].text_frame
content.text = "Key Routes"
for point in ["GET /api/students/ - List all students", "POST /api/students/ - Create student", "GET /api/recommendations/<id>/ - Get recommendations", "PUT /api/students/<id>/ - Update student", "DELETE /api/students/<id>/ - Delete student", "POST /api/auth/login/ - User authentication"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 6: Database & Auth Integration
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "2. Database & Auth Integration"
content = slide.placeholders[1].text_frame
content.text = "Secure Data Management"
for point in ["Successful database connection established", "JWT/Firebase Authentication implemented", "Secure password hashing (bcrypt/PBKDF2)", "Token-based authorization", "Protected routes for authenticated users", "Session management"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 7: Database Schema
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Database Schema"
content = slide.placeholders[1].text_frame
content.text = "Data Models"
for point in ["Student Model - Profile information", "Subject Model - Course details", "Recommendation Model - Learning resources", "User Model - Authentication data", "Proper relationships and foreign keys", "Migrations successfully applied"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 8: Full-Stack CRUD
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "3. Full-Stack CRUD Operations"
content = slide.placeholders[1].text_frame
content.text = "Frontend-Backend Integration"
for point in ["UI components fetch data from backend APIs", "Display student recommendations dynamically", "Create new student profiles with forms", "Update existing student data", "Delete operations with confirmation dialogs", "Real-time data synchronization"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 9: CRUD Implementation Details
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "CRUD Implementation"
content = slide.placeholders[1].text_frame
content.text = "Create"
for point in ["Form validation on frontend and backend", "Success/error notifications"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
p = content.add_paragraph()
p.text = "Read"
p.level = 0
for point in ["Dynamic data rendering", "Search and filter functionality"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
p = content.add_paragraph()
p.text = "Update & Delete"
p.level = 0
for point in ["Inline editing capabilities", "Confirmation before deletion"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 10: State Management
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "4. State Management"
content = slide.placeholders[1].text_frame
content.text = "Efficient Data Flow"
for point in ["React Hooks/Redux/Zustand implementation", "Centralized state management", "Smooth data flow between components", "Optimized re-rendering", "Persistent state handling", "Context API for global state"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 11: State Management Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "State Management Architecture"
content = slide.placeholders[1].text_frame
content.text = "Implementation"
for point in ["useState for local component state", "useEffect for side effects and API calls", "Custom hooks for reusable logic", "Global state for user authentication", "Loading and error states handled", "Optimistic UI updates"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 12: Error Handling & Security
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "5. Error Handling & Security"
content = slide.placeholders[1].text_frame
content.text = "Robust Security Implementation"
for point in ["Server-side input validation", "Secure HTTP headers (CORS, CSP, X-Frame-Options)", "Error logging mechanisms", "SQL injection prevention", "XSS protection", "CSRF token implementation"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 13: Security Measures
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Security Best Practices"
content = slide.placeholders[1].text_frame
content.text = "Implementation Details"
for point in ["Environment variables for sensitive data", "Rate limiting on API endpoints", "Input sanitization and validation", "Secure password storage", "HTTPS enforcement", "Regular security audits"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 14: Testing
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Testing & Quality Assurance"
content = slide.placeholders[1].text_frame
content.text = "Testing Strategy"
for point in ["API testing with Postman/Thunder Client", "Unit tests for backend logic", "Integration tests for API endpoints", "Frontend component testing", "Manual testing of user flows", "Bug tracking and resolution"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 15: Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Key Features"
content = slide.placeholders[1].text_frame
content.text = "System Capabilities"
for point in ["Personalized learning recommendations", "Multi-subject support (50+ subjects)", "Resource categorization (Books, Online, Practice)", "User authentication and profiles", "Responsive design for all devices", "Search and filter functionality"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 16: Challenges & Solutions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Challenges & Solutions"
content = slide.placeholders[1].text_frame
content.text = "Challenge: API Integration"
for point in ["Solution: Implemented proper error handling and retry logic"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
p = content.add_paragraph()
p.text = "Challenge: State Synchronization"
p.level = 0
for point in ["Solution: Used centralized state management"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
p = content.add_paragraph()
p.text = "Challenge: Security Vulnerabilities"
p.level = 0
for point in ["Solution: Implemented comprehensive security measures"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 17: Future Enhancements
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Enhancements"
content = slide.placeholders[1].text_frame
content.text = "Planned Features"
for point in ["Machine learning for better recommendations", "Real-time collaboration features", "Mobile application development", "Advanced analytics dashboard", "Integration with learning platforms", "Gamification elements"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 1

# Slide 18: Demo
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
tf = txBox.text_frame
tf.text = "Live Demo"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Slide 19: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = txBox.text_frame
tf.text = "Thank You"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf2 = txBox2.text_frame
tf2.text = "Questions?"
p2 = tf2.paragraphs[0]
p2.font.size = Pt(32)
p2.alignment = PP_ALIGN.CENTER

prs.save('Student_Recommendation_Phase2_Evaluation.pptx')
print("Presentation created: Student_Recommendation_Phase2_Evaluation.pptx")
